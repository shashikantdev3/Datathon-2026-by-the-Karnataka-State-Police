"""Populate the official KSP Datathon 2026 Prototype Submission Template in place."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TMPL = os.environ["TMPL"]; OUT = os.environ["OUT"]; IMG = os.environ["IMG"]
NAVY=RGBColor(0x0b,0x3d,0x91); BLUE=RGBColor(0x1d,0x4e,0xd8); INK=RGBColor(0x0f,0x17,0x2a)
GREY=RGBColor(0x47,0x55,0x69); RED=RGBColor(0xb9,0x1c,0x1c); GREEN=RGBColor(0x15,0x80,0x3d)

p = Presentation(TMPL)
def slide(i): return p.slides[i]

def add_text(s, top, lines, left=0.45, width=9.1, height=None, size=14, gap=6):
    h = height or (5.2-top)
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        txt, opts = (item, {}) if isinstance(item, str) else item
        para = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        para.space_after = Pt(opts.get("gap", gap)); para.space_before = Pt(0)
        para.alignment = PP_ALIGN.LEFT
        if opts.get("bullet"):
            txt = "•  " + txt
        run = para.add_run(); run.text = txt
        f = run.font
        f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", False)
        f.color.rgb = opts.get("color", INK); f.name = "Calibri"
    return tb

def add_img(s, path, left, top, width=None, height=None):
    kw = {}
    if width: kw["width"]=Inches(width)
    if height: kw["height"]=Inches(height)
    return s.shapes.add_picture(path, Inches(left), Inches(top), **kw)

# ---------- Slide 1: Team Details ----------
s = slide(0)
for sh in s.shapes:
    if sh.has_text_frame:
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip().startswith("Team name"):
                    run.text = "Team name: CrimeSense by DevWithData"

# ---------- Slide 2: Brief ----------
add_text(slide(1), 1.55, [
    ("CrimeSense is an Intelligent Conversational AI & Crime Analytics platform for the KSP crime database.", {"size":15,"bold":True,"color":NAVY,"gap":10}),
    ("Investigators, analysts and policymakers ask questions in plain English or Kannada (text or voice) and get instant answers — backed by a transparent evidence trail showing the exact data and reasoning used.", {"size":13,"gap":8}),
    ("Beyond simple retrieval, it discovers hidden links between crimes, offenders, victims, locations and money flows — powering crime-pattern discovery, criminal-network analysis, offender risk-scoring and proactive, preventive policing intelligence.", {"size":13,"gap":8}),
    ("Built for solo-team delivery and deployed on Catalyst by Zoho. Runs on fully synthetic, NCRB/KSP-structured data — no real personal data.", {"size":12,"color":GREY}),
])

# ---------- Slide 3: Opportunities / USP ----------
add_text(slide(2), 2.45, [
    ("How it is different:", {"size":13,"bold":True,"color":NAVY}),
    ("Not a dashboard or a search box — a conversational, context-aware analyst that explains every answer (SQL + reasoning), in English and Kannada.", {"size":12,"bullet":True}),
    ("How it solves the problem:", {"size":13,"bold":True,"color":NAVY,"gap":2}),
    ("One natural-language interface over FIRs, accused, victims, locations and transactions, plus network analysis, risk-scoring and money-trail tracing for real investigative leads.", {"size":12,"bullet":True}),
    ("USP:", {"size":13,"bold":True,"color":NAVY,"gap":2}),
    ("Explainable-by-design + bilingual + voice + offline NL→SQL (no API key needed) + role-based access & audit logs — purpose-built for law-enforcement accountability.", {"size":12,"bullet":True}),
])

# ---------- Slide 4: Features ----------
s = slide(3)
left_feats = [
    ("1. Conversational Crime Intelligence", "NL→SQL chat, EN/Kannada, voice, PDF export"),
    ("2. Criminal Network Analysis", "Co-offending graph, gangs, kingpin detection"),
    ("3. Crime Pattern & Trend Analytics", "Hotspots, yearly trend, seasonality"),
    ("4. Sociological Insights", "Age, gender, education, occupation patterns"),
    ("5. Offender Profiling & Risk Scoring", "Repeat offenders, transparent risk score"),
]
right_feats = [
    ("6. Investigator Decision Support", "Case summaries, status, similar cases"),
    ("7. Financial Crime & Money-Trail", "Multi-hop tracing, suspicious-txn flags"),
    ("8. Crime Forecasting", "Trend/seasonality, QuickML-ready"),
    ("9. Explainable AI", "Evidence trail: data refs + SQL + reasoning"),
    ("10. Secure Role-Based Access", "Roles + audit logs + data protection"),
]
def feat_lines(fs):
    out=[]
    for t,d in fs:
        out.append((t,{"size":12.5,"bold":True,"color":NAVY,"gap":1}))
        out.append((d,{"size":11,"color":GREY,"gap":7}))
    return out
add_text(s, 1.55, feat_lines(left_feats), left=0.45, width=4.55)
add_text(s, 1.55, feat_lines(right_feats), left=5.1, width=4.5)

# ---------- Slide 5: Process flow ----------
add_img(slide(4), os.path.join(IMG,"diagram_process.png"), left=0.4, top=1.6, width=9.2)

# ---------- Slide 6: Wireframe / mock ----------
add_img(slide(5), os.path.join(IMG,"01_chat.png"), left=2.51, top=1.5, height=3.4)
add_text(slide(5), 5.0, [("Conversational AI screen — question, answer, explainable evidence trail and auto-generated chart.",{"size":11,"color":GREY})], left=1.0, width=8.0, height=0.45)

# ---------- Slide 7: Architecture ----------
add_img(slide(6), os.path.join(IMG,"diagram_architecture.png"), left=1.55, top=1.6, width=6.9)

# ---------- Slide 8: Technologies ----------
add_text(slide(7), 1.55, [
    ("Frontend / App", {"size":13,"bold":True,"color":NAVY,"gap":2}),
    ("Python · Streamlit (interactive web UI) · matplotlib charts", {"size":12,"bullet":True,"gap":8}),
    ("AI / Analytics", {"size":13,"bold":True,"color":NAVY,"gap":2}),
    ("Explainable rule/slot NL→SQL engine · networkx (graph & centrality) · pandas · risk-scoring model · QuickML (forecasting)", {"size":12,"bullet":True,"gap":8}),
    ("Data", {"size":13,"bold":True,"color":NAVY,"gap":2}),
    ("SQLite (prototype) → Catalyst Data Store (serverless RDB + OLAP) · synthetic NCRB/KSP-structured dataset", {"size":12,"bullet":True,"gap":8}),
    ("Languages / Other", {"size":13,"bold":True,"color":NAVY,"gap":2}),
    ("English + Kannada (i18n) · Web Speech API (voice) · reportlab (PDF export) · Faker (data generation)", {"size":12,"bullet":True}),
])

# ---------- Slide 9: Catalyst services ----------
add_text(slide(8), 1.55, [
    ("AppSail", {"size":13,"bold":True,"color":GREEN,"gap":1}),
    ("Serverless Python web hosting for the Streamlit app (the deployed link).", {"size":12,"bullet":True,"gap":7}),
    ("Data Store", {"size":13,"bold":True,"color":GREEN,"gap":1}),
    ("Serverless relational DB with built-in OLAP for crime records & analytical queries.", {"size":12,"bullet":True,"gap":7}),
    ("Catalyst Functions (FaaS)", {"size":13,"bold":True,"color":GREEN,"gap":1}),
    ("NL-query microservice and PDF-export function.", {"size":12,"bullet":True,"gap":7}),
    ("QuickML", {"size":13,"bold":True,"color":GREEN,"gap":1}),
    ("No-code ML pipelines for crime forecasting and offender risk-scoring.", {"size":12,"bullet":True,"gap":7}),
    ("Authentication & Stratus", {"size":13,"bold":True,"color":GREEN,"gap":1}),
    ("Role-based access (investigator/analyst/supervisor/policymaker) + storage for exported PDFs.", {"size":12,"bullet":True}),
])

# ---------- Slide 10: Estimated cost ----------
add_text(slide(9), 1.55, [
    ("Prototype runs on Catalyst's free developer tier — effectively zero cost.", {"size":13,"bold":True,"color":NAVY,"gap":10}),
    ("AppSail hosting (1 app, 512 MB):  Free tier → ~₹0–1,500 / month at pilot scale", {"size":12,"bullet":True,"gap":6}),
    ("Data Store (crime records):  Free tier → pay-as-you-grow", {"size":12,"bullet":True,"gap":6}),
    ("Catalyst Functions + QuickML:  Free quota covers prototype workloads", {"size":12,"bullet":True,"gap":6}),
    ("No third-party LLM API cost — NL→SQL engine runs fully offline.", {"size":12,"bullet":True,"gap":6}),
    ("Indicative production pilot (single district): well under ₹10,000 / month.", {"size":12,"color":GREY}),
])

# ---------- Slide 11: Snapshots ----------
s = slide(10)
shots = ["02_patterns.png","03_network.png","04_offenders.png","05_financial.png"]
poss = [(0.35,1.5),(5.05,1.5),(0.35,3.45),(5.05,3.45)]
for img,(lx,ty) in zip(shots,poss):
    add_img(s, os.path.join(IMG,img), left=lx, top=ty, width=4.55)

# ---------- Slide 12: Performance / Benchmark ----------
add_text(slide(11), 1.55, [
    ("Prototype benchmark (synthetic KSP-style dataset)", {"size":13,"bold":True,"color":NAVY,"gap":10}),
    ("Dataset: 1,500 FIRs · 600 accused · 1,571 victims · 669 financial transactions · 4,219 case-log events across 12 Karnataka districts.", {"size":12,"bullet":True,"gap":6}),
    ("NL query response: < 200 ms (offline rule/slot engine + indexed SQLite).", {"size":12,"bullet":True,"gap":6}),
    ("Intent coverage: count, trend, hotspots, breakdown, offenders, gangs, status, financial, case-detail — 9 intents, context-aware follow-ups.", {"size":12,"bullet":True,"gap":6}),
    ("Network analysis: 117 linked offenders · 168 relationships · 14 organised-crime cells detected via modularity.", {"size":12,"bullet":True,"gap":6}),
    ("Explainability: 100% of answers return SQL + reasoning trail. Read-only — zero write/DDL executed.", {"size":12,"bullet":True}),
])

# ---------- Slide 13: Links ----------
add_text(slide(12), 2.4, [
    ("GitHub Public Repository", {"size":13,"bold":True,"color":NAVY,"gap":1}),
    ("https://github.com/<your-username>/Datathon-2026-KSP-CrimeSense   ← replace with your public repo URL", {"size":12,"color":GREY,"gap":12}),
    ("Demo Video Link (3 minutes)", {"size":13,"bold":True,"color":NAVY,"gap":1}),
    ("https://youtu.be/<unlisted-id>   or   public Google Drive link   ← add after recording", {"size":12,"color":GREY,"gap":12}),
    ("Deployed Link (on Catalyst by Zoho)", {"size":13,"bold":True,"color":NAVY,"gap":1}),
    ("https://crimesense-<id>.catalystserverless.com   ← add after `catalyst deploy`", {"size":12,"color":GREY,"gap":12}),
    ("Note: all three links must be public & working before final submission (test in incognito).", {"size":11,"color":RED}),
])

# ---------- Slide 14: Future development ----------
add_text(slide(13), 1.55, [
    ("LLM-augmented NL engine (Catalyst AI / Zia) for open-ended questions, while keeping the explainable offline path as fallback.", {"size":12,"bullet":True,"gap":7}),
    ("Live integration with KSP CCTNS / FIR systems and real (access-controlled) data, with full audit & data-protection compliance.", {"size":12,"bullet":True,"gap":7}),
    ("Geospatial hotspot maps and QuickML forecasting for early-warning alerts on repeat/gang crime.", {"size":12,"bullet":True,"gap":7}),
    ("Full Kannada voice assistant (STT + TTS) for field officers.", {"size":12,"bullet":True,"gap":7}),
    ("Bias & fairness auditing of risk scores, with human-in-the-loop review workflows.", {"size":12,"bullet":True}),
])

p.save(OUT)
print("saved", OUT, "slides", len(p.slides))
